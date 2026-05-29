# ============================================================
# rag_engine.py — Document RAG Pipeline for StudyPilot
#
# IMPORTANT DESIGN DECISION:
# Your project already uses FAISS + sentence-transformers in memory.py
# So this file uses the SAME stack instead of adding ChromaDB.
# This means: one embedding model loaded once, one consistent vector approach.
#
# NEW packages needed (add to requirements.txt):
#   pymupdf        → read PDFs  (imported as fitz)
#   python-docx    → read DOCX files
#   python-pptx    → read PPTX files
# ============================================================

import os
import pickle
import uuid

import numpy as np
import faiss
import fitz                          # pip install pymupdf
from docx import Document as DocxDoc # pip install python-docx
from pptx import Presentation        # pip install python-pptx
from sentence_transformers import SentenceTransformer
from datetime import datetime
from typing import List, Dict
import re

# ============================================================
# CONSTANTS
# ============================================================

EMBEDDING_DIM    = 384
RAG_INDEX_PATH   = "rag_faiss_index.bin"    # Separate file from memory.py's index
RAG_META_PATH    = "rag_faiss_metadata.pkl"  # So document chunks never mix with chat history

# ============================================================
# SHARED EMBEDDING MODEL
# ============================================================
# We load the same model as memory.py ("all-MiniLM-L6-v2").
# Python caches imported modules, so if memory.py is already loaded,
# this doesn't download or initialise a second model — it reuses.
# However to be safe and explicit we create our own reference here.

_embedder = SentenceTransformer("all-MiniLM-L6-v2")


# ============================================================
# FAISS INDEX MANAGEMENT (mirrors memory.py's pattern)
# ============================================================

def _load_rag_index():
    """
    Loads the RAG FAISS index and metadata from disk if they exist.
    If not (first run), creates a fresh empty index.

    faiss.IndexFlatIP — 'IP' = Inner Product similarity.
    Since our embeddings are L2-normalised, inner product == cosine similarity.
    This is exactly what memory.py uses too, so behaviour is consistent.
    """
    if os.path.exists(RAG_INDEX_PATH) and os.path.exists(RAG_META_PATH):
        index = faiss.read_index(RAG_INDEX_PATH)
        with open(RAG_META_PATH, "rb") as f:
            metadata = pickle.load(f)
    else:
        index    = faiss.IndexFlatIP(EMBEDDING_DIM)
        metadata = []
    return index, metadata


def _save_rag_index(index, metadata):
    """Saves the FAISS index and metadata to disk."""
    faiss.write_index(index, RAG_INDEX_PATH)
    with open(RAG_META_PATH, "wb") as f:
        pickle.dump(metadata, f)


# Load once at module level — same pattern as memory.py
_rag_index, _rag_metadata = _load_rag_index()


def _embed(text: str) -> np.ndarray:
    """
    Converts text to a normalised 384-dim vector.
    normalize_embeddings=True ensures inner product == cosine similarity.
    Returns shape (1, 384) float32 — required by faiss.add().
    """
    vector = _embedder.encode(text, normalize_embeddings=True)
    return vector.astype("float32").reshape(1, -1)


# ============================================================
# TEXT EXTRACTION
# ============================================================

def extract_text_from_pdf(pdf_path: str) -> str:
    """
    Reads every page of a PDF and returns all text as one string.

    fitz.open()       — opens the PDF file
    doc.load_page(i)  — loads page at index i (0-based)
    page.get_text()   — extracts plain text from that page
    """
    doc  = fitz.open(pdf_path)
    text = ""
    for i in range(len(doc)):
        page  = doc.load_page(i)
        text += page.get_text()
    doc.close()
    # Collapse multiple whitespace/newlines into single spaces
    return re.sub(r'\s+', ' ', text).strip()


def extract_text_from_docx(docx_path: str) -> str:
    """
    Reads a .docx file paragraph by paragraph and returns all text.

    DocxDoc(path)      — opens the Word document
    doc.paragraphs     — list of paragraph objects
    para.text          — the plain text of one paragraph
    """
    doc  = DocxDoc(docx_path)
    text = " ".join(para.text for para in doc.paragraphs if para.text.strip())
    return re.sub(r'\s+', ' ', text).strip()


def extract_text_from_pptx(pptx_path: str) -> str:
    """
    Reads a .pptx file slide by slide and returns all text.

    Presentation(path)       — opens the PowerPoint file
    prs.slides               — list of slide objects
    slide.shapes             — list of shapes (text boxes, images, etc.)
    shape.has_text_frame     — True if shape contains text
    shape.text_frame.text    — all text in that shape
    """
    prs  = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text_frame.text + " "
    return re.sub(r'\s+', ' ', text).strip()


def extract_text_from_txt(txt_path: str) -> str:
    """Reads a plain .txt file and returns its content."""
    with open(txt_path, "r", encoding="utf-8", errors="ignore") as f:
        return re.sub(r'\s+', ' ', f.read()).strip()


def extract_text(file_path: str) -> str:
    """
    Router function — detects file type by extension and calls
    the correct extractor. Raises ValueError for unsupported types.
    """
    ext = os.path.splitext(file_path)[1].lower()
    # os.path.splitext("notes.pdf") → ("notes", ".pdf")
    # .lower() makes the check case-insensitive (.PDF == .pdf)

    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext == ".txt":
        return extract_text_from_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}. Supported: PDF, DOCX, PPTX, TXT")


# ============================================================
# CHUNKING
# ============================================================

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
    """
    Splits a long string into overlapping chunks.

    Why chunking:
    - LLMs have token limits; a 200-page book can't go into one prompt.
    - FAISS searches work better on focused, small pieces.
    - Overlap prevents losing meaning at chunk boundaries.

    Example with chunk_size=500, overlap=50:
    - Chunk 1: chars 0   → 500
    - Chunk 2: chars 450 → 950   (50 char overlap with chunk 1)
    - Chunk 3: chars 900 → 1400
    """
    chunks = []
    start  = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks


# ============================================================
# STORE DOCUMENT
# ============================================================

def store_document(file_path: str, user_id: str, doc_name: str) -> int:
    """
    Full pipeline: file → text → chunks → embeddings → FAISS storage.

    Each chunk is stored as a row in the FAISS index with metadata so we can:
    - Filter by user_id (user isolation)
    - Filter by doc_name (per-document search)
    - Return the source name in results

    Args:
        file_path: Absolute path to the uploaded file on disk
        user_id:   Firebase UID — keeps each user's docs isolated
        doc_name:  Friendly name shown in the UI (e.g. "Physics Notes")

    Returns:
        Number of chunks stored
    """
    global _rag_index, _rag_metadata

    # Step 1: Extract text based on file type
    text = extract_text(file_path)
    if not text.strip():
        raise ValueError("No text could be extracted. The file may be empty or image-only.")

    # Step 2: Split into chunks
    chunks = chunk_text(text)

    # Step 3: Embed all chunks at once (batch is faster than one-by-one)
    # _embedder.encode(list) → numpy array of shape (num_chunks, 384)
    vectors = _embedder.encode(chunks, normalize_embeddings=True).astype("float32")

    # Step 4: Remove old chunks from this document if re-uploading
    # This prevents duplicate chunks if the user uploads the same doc twice
    _rag_metadata = [
        m for m in _rag_metadata
        if not (m["user_id"] == user_id and m["doc_name"] == doc_name)
    ]
    # Rebuild index without the deleted entries
    # (FAISS IndexFlatIP doesn't support in-place deletion, so we rebuild)
    _rag_index = faiss.IndexFlatIP(EMBEDDING_DIM)
    if _rag_metadata:
        all_vectors = np.array(
            [m["vector"] for m in _rag_metadata], dtype="float32"
        )
        _rag_index.add(all_vectors)

    # Step 5: Add new chunks to the index
    _rag_index.add(vectors)

    # Step 6: Add metadata for each chunk (parallel to the FAISS rows)
    for i, (chunk, vector) in enumerate(zip(chunks, vectors)):
        _rag_metadata.append({
            "id":        str(uuid.uuid4()),
            "user_id":   user_id,
            "doc_name":  doc_name,
            "chunk_idx": i,
            "content":   chunk,
            "vector":    vector,   # Store vector so we can rebuild index on deletion
            "timestamp": datetime.now().isoformat()
        })

    # Step 7: Persist to disk
    _save_rag_index(_rag_index, _rag_metadata)

    return len(chunks)


# ============================================================
# RETRIEVE RELEVANT CHUNKS
# ============================================================

def retrieve_relevant_chunks(query: str, user_id: str, top_k: int = 5) -> List[Dict]:
    """
    Finds the top_k most semantically relevant chunks for a query.

    How it works:
    1. Embed the query into a vector
    2. FAISS searches all stored vectors for the most similar ones
       (inner product with normalised vectors = cosine similarity)
    3. Filter results to only this user's documents
    4. Return chunk text, source name, and similarity score

    Args:
        query:   What to search for (e.g., "Newton's laws", "photosynthesis")
        user_id: Only search this user's documents
        top_k:   Max number of chunks to return

    Returns:
        List of dicts with keys: text, source, score
    """
    global _rag_index, _rag_metadata

    if _rag_index.ntotal == 0:
        return []

    # Embed the query
    query_vector = _embed(query)

    # Search — returns distances (scores) and indices
    # We ask for more than top_k because some results may belong to other users
    # and get filtered out below
    k = min(top_k * 4, _rag_index.ntotal)
    scores, indices = _rag_index.search(query_vector, k)

    results = []
    for idx, score in zip(indices[0], scores[0]):
        if idx == -1:
            continue
        entry = _rag_metadata[idx]
        # Only include chunks belonging to this user
        if entry["user_id"] != user_id:
            continue
        results.append({
            "text":   entry["content"],
            "source": entry["doc_name"],
            "score":  float(score)
        })
        if len(results) == top_k:
            break

    return results


# ============================================================
# LIST & DELETE (for Streamlit sidebar management)
# ============================================================

def list_user_documents(user_id: str) -> List[str]:
    """Returns list of unique document names the user has indexed."""
    seen  = set()
    names = []
    for m in _rag_metadata:
        if m["user_id"] == user_id and m["doc_name"] not in seen:
            seen.add(m["doc_name"])
            names.append(m["doc_name"])
    return names


def delete_user_document(user_id: str, doc_name: str) -> bool:
    """
    Removes all chunks belonging to a specific document and rebuilds the index.
    Returns True on success.
    """
    global _rag_index, _rag_metadata

    before = len(_rag_metadata)
    _rag_metadata = [
        m for m in _rag_metadata
        if not (m["user_id"] == user_id and m["doc_name"] == doc_name)
    ]

    if len(_rag_metadata) == before:
        return False  # Nothing was deleted

    # Rebuild FAISS index without the deleted document's vectors
    _rag_index = faiss.IndexFlatIP(EMBEDDING_DIM)
    if _rag_metadata:
        all_vectors = np.array(
            [m["vector"] for m in _rag_metadata], dtype="float32"
        )
        _rag_index.add(all_vectors)

    _save_rag_index(_rag_index, _rag_metadata)
    return True


def has_documents(user_id: str) -> bool:
    """Quick check: does this user have any indexed documents?"""
    return any(m["user_id"] == user_id for m in _rag_metadata)
